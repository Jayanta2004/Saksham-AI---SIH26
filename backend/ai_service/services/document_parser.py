import os
import re
from typing import List, Dict, Any, Optional

class DocumentParser:
    """
    Document Ingestion & Parsing Service for PDFs, PPT/PPTXs, and Text files.
    Extracts text, slide titles, tables, section metadata, and chunks text semantically.
    """

    @staticmethod
    def parse_raw_text(text: str, source_name: str = "Raw Text", chunk_size: int = 500, overlap: int = 50) -> List[Dict[str, Any]]:
        """
        Parses raw text directly into semantic chunks with metadata.
        """
        clean_text = DocumentParser._clean_text(text)
        if not clean_text:
            return []
        items = [{"text": clean_text, "page_number": 1, "type": "raw_text"}]
        return DocumentParser.chunk_document(items, chunk_size=chunk_size, chunk_overlap=overlap)

    @staticmethod
    def parse_pdf(file_path: str) -> List[Dict[str, Any]]:
        """
        Extract pages and content from a PDF file using pypdf or PyPDF2 fallback.
        """
        pages_content = []
        try:
            try:
                import pypdf
                reader = pypdf.PdfReader(file_path)
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    clean_text = DocumentParser._clean_text(text)
                    if clean_text:
                        pages_content.append({
                            "page_number": i + 1,
                            "text": clean_text,
                            "type": "page"
                        })
            except ImportError:
                import PyPDF2
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for i, page in enumerate(reader.pages):
                        text = page.extract_text() or ""
                        clean_text = DocumentParser._clean_text(text)
                        if clean_text:
                            pages_content.append({
                                "page_number": i + 1,
                                "text": clean_text,
                                "type": "page"
                            })
        except Exception as e:
            # Fallback if binary read fails or is plain text
            print(f"[DocumentParser] PDF parsing note: {e}. Attempting fallback reading.")
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    raw = f.read()
                    pages_content.append({
                        "page_number": 1,
                        "text": DocumentParser._clean_text(raw),
                        "type": "text_fallback"
                    })
            except Exception as read_err:
                print(f"[DocumentParser] Error reading file: {read_err}")
                
        return pages_content

    @staticmethod
    def parse_pptx(file_path: str) -> List[Dict[str, Any]]:
        """
        Extract slides, titles, shapes, and presenter notes from PowerPoint presentations.
        """
        slides_content = []
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                title = f"Slide {i + 1}"
                
                # Check for slide title
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            p_text = paragraph.text.strip()
                            if p_text and p_text != title:
                                slide_texts.append(p_text)
                    elif shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_vals = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_vals:
                                slide_texts.append(" | ".join(row_vals))

                # Slide notes if any
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_texts.append(f"[Presenter Notes: {notes}]")

                full_slide_text = f"{title}\n" + "\n".join(slide_texts)
                clean_text = DocumentParser._clean_text(full_slide_text)
                if clean_text:
                    slides_content.append({
                        "slide_number": i + 1,
                        "title": title,
                        "text": clean_text,
                        "type": "slide"
                    })
        except Exception as e:
            print(f"[DocumentParser] PPTX parsing error: {e}. Using text fallback.")
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    raw = f.read()
                    slides_content.append({
                        "slide_number": 1,
                        "title": "Presentation Overview",
                        "text": DocumentParser._clean_text(raw),
                        "type": "text_fallback"
                    })
            except Exception as read_err:
                print(f"[DocumentParser] Error reading PPTX file: {read_err}")
                
        return slides_content

    @staticmethod
    def parse_text_or_markdown(file_path: str) -> List[Dict[str, Any]]:
        """
        Parse raw text or markdown guidelines.
        """
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                return [{
                    "page_number": 1,
                    "text": DocumentParser._clean_text(content),
                    "type": "text"
                }]
        except Exception as e:
            print(f"[DocumentParser] Error reading text file: {e}")
            return []

    @staticmethod
    def parse_file(file_path: str) -> List[Dict[str, Any]]:
        """
        Dispatch parser based on file extension.
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext in ['.pdf']:
            return DocumentParser.parse_pdf(file_path)
        elif ext in ['.ppt', '.pptx']:
            return DocumentParser.parse_pptx(file_path)
        else:
            return DocumentParser.parse_text_or_markdown(file_path)

    @staticmethod
    def chunk_document(
        parsed_items: List[Dict[str, Any]], 
        chunk_size: int = 600, 
        chunk_overlap: int = 120
    ) -> List[Dict[str, Any]]:
        """
        Splits document contents into semantic overlapping chunks with rich metadata.
        """
        chunks = []
        chunk_counter = 1

        for item in parsed_items:
            text = item.get("text", "")
            meta_label = f"Page {item.get('page_number')}" if "page_number" in item else f"Slide {item.get('slide_number')} ({item.get('title', '')})"
            
            # Sentence-aware chunking
            sentences = re.split(r'(?<=[.?!])\s+', text)
            current_chunk = []
            current_len = 0

            for sentence in sentences:
                s_len = len(sentence.split())
                if current_len + s_len > chunk_size and current_chunk:
                    chunk_text = " ".join(current_chunk)
                    chunks.append({
                        "chunk_id": f"chk_{chunk_counter:04d}",
                        "text": chunk_text,
                        "source_location": meta_label,
                        "token_count": len(chunk_text.split()),
                        "metadata": item
                    })
                    chunk_counter += 1
                    
                    # Keep overlap words
                    overlap_words = " ".join(current_chunk).split()[-chunk_overlap:]
                    current_chunk = [" ".join(overlap_words), sentence]
                    current_len = len(overlap_words) + s_len
                else:
                    current_chunk.append(sentence)
                    current_len += s_len

            if current_chunk:
                chunk_text = " ".join(current_chunk)
                chunks.append({
                    "chunk_id": f"chk_{chunk_counter:04d}",
                    "text": chunk_text,
                    "source_location": meta_label,
                    "token_count": len(chunk_text.split()),
                    "metadata": item
                })
                chunk_counter += 1

        return chunks

    @staticmethod
    def _clean_text(text: str) -> str:
        if not text:
            return ""
        # Remove null bytes, multiple spaces, repeated newlines
        cleaned = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\xff]', '', text)
        cleaned = re.sub(r'\r\n|\r', '\n', cleaned)
        cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
        cleaned = re.sub(r'[ \t]+', ' ', cleaned)
        return cleaned.strip()
